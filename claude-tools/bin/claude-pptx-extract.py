#!/usr/bin/env python3
"""Extract text, speaker notes, and slide images from a .pptx file."""

import sys
import io
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

DPI = 150
EMU_PER_INCH = 914400


def emu_to_px(emu):
    """Convert EMU (English Metric Units) to pixels at the configured DPI."""
    return int(emu * DPI / EMU_PER_INCH)


def load_font(size):
    """Try to load a system font at the given size, falling back to default."""
    for family in ("Helvetica", "DejaVuSans", "DejaVu Sans", "Arial"):
        try:
            return ImageFont.truetype(family, size)
        except (OSError, IOError):
            continue
    # PIL default bitmap font (ignores size, but always works)
    return ImageFont.load_default()


def get_slide_title(slide, index):
    """Return the slide title text or a fallback like 'Slide N'."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()
        if title:
            return title
    return f"Slide {index}"


def get_slide_text(slide):
    """Extract all text from shapes with text frames, excluding the title."""
    lines = []
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
    return "\n".join(lines)


def get_speaker_notes(slide):
    """Extract speaker notes text from the notes slide, if present."""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            text = notes_slide.notes_text_frame.text.strip()
            if text:
                return text
    return ""


def get_run_color(runs):
    """Extract hex color string from the first run with a color, or None."""
    for run in runs:
        try:
            rgb = run.font.color.rgb
            if rgb is not None:
                return f"#{rgb}"
        except (AttributeError, TypeError):
            continue
    return None


def draw_text_wrapped(draw, text, x, y, max_width, font, fill="black"):
    """Draw text with word-wrapping within a bounding width."""
    words = text.split()
    if not words:
        return y

    lines = []
    current_line = words[0]
    for word in words[1:]:
        test_line = current_line + " " + word
        bbox = draw.textbbox((0, 0), test_line, font=font)
        line_width = bbox[2] - bbox[0]
        if line_width <= max_width:
            current_line = test_line
        else:
            lines.append(current_line)
            current_line = word
    lines.append(current_line)

    for line in lines:
        draw.text((x, y), line, fill=fill, font=font)
        bbox = draw.textbbox((0, 0), line, font=font)
        line_height = bbox[3] - bbox[1]
        y += line_height + 4  # small line spacing
    return y


def get_font_size_from_shape(shape):
    """Get the most common font size from a shape's text runs, in points."""
    sizes = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    sizes.append(run.font.size.pt)
    if sizes:
        return max(int(max(set(sizes), key=sizes.count)), 6)
    return 14  # default


def render_slide(prs, slide, canvas_width, canvas_height):
    """Render a slide to a PIL Image."""
    img = Image.new("RGB", (canvas_width, canvas_height), "white")
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        try:
            left = emu_to_px(shape.left) if shape.left is not None else 0
            top = emu_to_px(shape.top) if shape.top is not None else 0
            width = emu_to_px(shape.width) if shape.width is not None else 100
            height = emu_to_px(shape.height) if shape.height is not None else 100

            # --- Embedded images ---
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blob = shape.image.blob
                    embedded = Image.open(io.BytesIO(image_blob))
                    embedded = embedded.convert("RGB")
                    embedded = embedded.resize((width, height), Image.LANCZOS)
                    img.paste(embedded, (left, top))
                except Exception:
                    # Gray placeholder on failure
                    draw.rectangle(
                        [left, top, left + width, top + height],
                        fill="#CCCCCC",
                        outline="#999999",
                    )
                    placeholder_font = load_font(12)
                    draw.text(
                        (left + 4, top + 4), "[image]", fill="#666666", font=placeholder_font
                    )
                continue

            # --- Group shapes: try to recurse into pictures inside groups ---
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for child in shape.shapes:
                        cl = emu_to_px(child.left) if child.left is not None else 0
                        ct = emu_to_px(child.top) if child.top is not None else 0
                        cw = emu_to_px(child.width) if child.width is not None else 100
                        ch = emu_to_px(child.height) if child.height is not None else 100
                        if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image_blob = child.image.blob
                                embedded = Image.open(io.BytesIO(image_blob))
                                embedded = embedded.convert("RGB")
                                embedded = embedded.resize((cw, ch), Image.LANCZOS)
                                img.paste(embedded, (cl, ct))
                            except Exception:
                                draw.rectangle(
                                    [cl, ct, cl + cw, ct + ch],
                                    fill="#CCCCCC",
                                    outline="#999999",
                                )
                        elif child.has_text_frame:
                            font_size = 14
                            child_color = "black"
                            if hasattr(child, "text_frame"):
                                for para in child.text_frame.paragraphs:
                                    for run in para.runs:
                                        if run.font.size is not None:
                                            font_size = max(int(run.font.size.pt), 6)
                                        c = get_run_color([run])
                                        if c is not None:
                                            child_color = c
                                        break
                                    if font_size != 14:
                                        break
                            font = load_font(font_size)
                            text = child.text_frame.text.strip()
                            if text:
                                draw_text_wrapped(
                                    draw, text, cl + 4, ct + 4, cw - 8, font,
                                    fill=child_color,
                                )
                except Exception:
                    pass
                continue

            # --- Tables ---
            if shape.has_table:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                if rows == 0 or cols == 0:
                    continue

                cell_w = width // cols if cols > 0 else width
                cell_h = height // rows if rows > 0 else height
                table_font = load_font(11)

                for r_idx in range(rows):
                    for c_idx in range(cols):
                        cx = left + c_idx * cell_w
                        cy = top + r_idx * cell_h
                        # Draw cell border
                        draw.rectangle(
                            [cx, cy, cx + cell_w, cy + cell_h],
                            outline="#333333",
                        )
                        # Draw cell text
                        try:
                            cell = table.cell(r_idx, c_idx)
                            cell_text = cell.text.strip()
                            if cell_text:
                                draw_text_wrapped(
                                    draw, cell_text, cx + 4, cy + 4, cell_w - 8, table_font
                                )
                        except Exception:
                            pass
                continue

            # --- Text shapes ---
            if shape.has_text_frame:
                font_size = get_font_size_from_shape(shape)
                font = load_font(font_size)
                y_cursor = top + 4
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        color = get_run_color(para.runs) or "black"
                        y_cursor = draw_text_wrapped(
                            draw, text, left + 4, y_cursor, width - 8, font,
                            fill=color,
                        )
                    else:
                        y_cursor += font_size  # blank line spacing
        except Exception as e:
            print(f"  Warning: failed to render shape on slide: {e}", file=sys.stderr)

    return img


def main():
    if len(sys.argv) != 2:
        print(
            "Usage: claude-pptx-extract.py <presentation.pptx>",
            file=sys.stderr,
        )
        sys.exit(1)

    pptx_path = Path(sys.argv[1]).resolve()
    if not pptx_path.is_file():
        print(f"ERROR: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    if pptx_path.suffix.lower() != ".pptx":
        print(f"ERROR: Expected .pptx file, got: {pptx_path.suffix}", file=sys.stderr)
        sys.exit(1)

    stem = pptx_path.stem
    out_dir = pptx_path.parent

    prs = Presentation(str(pptx_path))

    # Canvas dimensions from slide layout
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    canvas_w = emu_to_px(slide_width_emu)
    canvas_h = emu_to_px(slide_height_emu)

    # Determine presentation title from first slide title
    slides = list(prs.slides)
    if not slides:
        print("ERROR: Presentation has no slides.", file=sys.stderr)
        sys.exit(1)

    pres_title = get_slide_title(slides[0], 1)

    md_lines = [f"# {pres_title}", ""]

    for i, slide in enumerate(slides, start=1):
        slide_num = f"{i:02d}"
        title = get_slide_title(slide, i)
        text = get_slide_text(slide)
        notes = get_speaker_notes(slide)

        # Render slide image
        img_filename = f"{stem}-slide-{slide_num}.png"
        img_path = out_dir / img_filename
        try:
            slide_img = render_slide(prs, slide, canvas_w, canvas_h)
            slide_img.save(str(img_path), "PNG")
            print(f"  {img_path}")
        except Exception as e:
            print(f"  Warning: could not render slide {i}: {e}", file=sys.stderr)

        # Build markdown section
        md_lines.append(f"## Slide {i}: {title}")
        md_lines.append("")
        md_lines.append(f"![Slide {i}](./{img_filename})")
        md_lines.append("")

        if text:
            md_lines.append(text)
            md_lines.append("")

        if notes:
            md_lines.append("**Speaker Notes:**")
            md_lines.append(notes)
            md_lines.append("")

    # Write content markdown
    md_path = out_dir / f"{stem}-content.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")
    print(f"  {md_path}")

    print(f"Done. {len(slides)} slide(s) processed.")


if __name__ == "__main__":
    main()
